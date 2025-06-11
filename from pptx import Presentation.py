from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

# Define positions
left_margin = Inches(0.5)
top_margin = Inches(0.5)
spacing_y = Inches(1.0)
spacing_x = Inches(2.5)

# Helper function to add a labeled shape
def add_box(slide, text, left, top, width=Inches(2), height=Inches(0.6), color=(91, 155, 213)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.text = text
    text_frame = shape.text_frame
    text_frame.paragraphs[0].font.size = Pt(12)
    text_frame.paragraphs[0].font.bold = True
    return shape

# Add site 3S2I components
slide.shapes.add_textbox(left_margin, top_margin - Inches(0.3), Inches(3), Inches(0.3)).text = "🏢 Site 3S2I"
add_box(slide, "WatchGuard T80", left_margin, top_margin)
add_box(slide, "Switch", left_margin, top_margin + spacing_y)
add_box(slide, "NAS", left_margin, top_margin + spacing_y * 2)
add_box(slide, "Serveur ERP", left_margin, top_margin + spacing_y * 3)
add_box(slide, "Serveur Redondance", left_margin, top_margin + spacing_y * 4)

# Add site BIOBULLE components
right = left_margin + spacing_x * 2.5
slide.shapes.add_textbox(right, top_margin - Inches(0.3), Inches(3), Inches(0.3)).text = "🏭 Site BIOBULLE"
add_box(slide, "WatchGuard", right, top_margin)
add_box(slide, "Postes utilisateurs", right, top_margin + spacing_y)

# Add BOVPN connector (simple line between the two firewalls)
line = slide.shapes.add_connector(
    connector_type=1,  # straight line
    begin_x=left_margin + Inches(2),
    begin_y=top_margin + Inches(0.3),
    end_x=right,
    end_y=top_margin + Inches(0.3)
)
line.line.color.rgb = RGBColor(0, 0, 0)
line.line.width = Pt(2)

# Label the BOVPN
textbox = slide.shapes.add_textbox((left_margin + right) / 2 - Inches(0.5), top_margin - Inches(0.1), Inches(1.5), Inches(0.3))
textbox.text_frame.text = "🔗 BOVPN"
textbox.text_frame.paragraphs[0].font.size = Pt(12)
textbox.text_frame.paragraphs[0].font.bold = True

# Save presentation
pptx_path = "/mnt/data/Schema_Reseau_3S2I_Biobulle.pptx"
prs.save(pptx_path)

pptx_path
